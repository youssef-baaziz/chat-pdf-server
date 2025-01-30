import os
import json
from dotenv import load_dotenv
import database
import chatpdf_processor
import aiofiles
from datetime import datetime
import random
import string
import asyncio
import fitz
from PIL import Image
import pytesseract
from io import BytesIO
import time
from callbacks import CustomStreamingCallbackHandler
from pptx import Presentation
                
# Load environment variables
load_dotenv()

# Initialize session state variables
session_state = {
    "history": {"id": [], "question": [], "answer": []},
    "hold_descriptions": [],
    "hold_files": [],
    "file_names": '',
    "conversation": [],
}

def initialize_data():
    session_state['hold_files'] = []
    session_state['conversation'] = []
    session_state['hold_descriptions'] = []
    session_state['file_names'] = ''
    session_state["history"] = {'id': [], 'question': [], 'answer': []}
    
word_json = ".json"
import zipfile
def extract_images_from_pptx(doc):
    pptx_zip = zipfile.ZipFile(doc, 'r')
    image_files = []

    for file_name in pptx_zip.namelist():
        if file_name.startswith('ppt/media/'):
            image_files.append(file_name)
    
    return pptx_zip, image_files

def get_text_from_image_pptx(image_bytes):
    text_image = ""
    try:
        # Create an image object from the extracted image bytes
        image = Image.open(BytesIO(image_bytes))

        # Use Tesseract to extract text from the image
        text = pytesseract.image_to_string(image)
        text_image += text + "\n"

    finally:
        image.close()
    
    return text_image

def get_text_from_image(images,page_num,doc):
    text_image = ""
    print(f"Page {page_num + 1} contains images.")
                        
    for img_index, img in enumerate(images):
        xref = img[0]  # The image XREF index
        base_image = doc.extract_image(xref)
        image_bytes = base_image["image"]  # Raw image data
        image_ext = base_image["ext"]  # Image file extension (e.g., 'png')

        try:
            # Create an image object from the extracted image bytes
            image = Image.open(BytesIO(image_bytes))

            # Save the image temporarily
            image_filename = f"extracted_image_page{page_num+1}_img{img_index+1}.{image_ext}"
            image.save(image_filename)
            
            # Use Tesseract to extract text from the image
            text = pytesseract.image_to_string(image)
            
            text_image += text + "\n"

        finally:
            image.close()

            if os.path.exists(image_filename):
                os.remove(image_filename)
                print(f"Deleted image: {image_filename}")
            else:
                print(f"Failed to delete")
                    
    return text_image             
                    
def get_text_from_docc(docs, descriptions=None):
    text = ""
    
    for i, doc in enumerate(docs):
        file_content = ""
        file_name = os.path.basename(doc)
        description = descriptions[i] if descriptions and i < len(descriptions) else "No description provided"
        
        file_content += f"Name of the file: {file_name}\n"
        file_content += f"Description: {description}\n"
        file_content += "Content:\n"
        pytesseract.pytesseract.tesseract_cmd = r'C:\Users\ybaaziz\AppData\Local\Programs\Tesseract-OCR\tesseract.exe'
        
        try:
            # Handle PDF files
            if doc.endswith(".pdf"):
                pdf_document = fitz.open(doc)
                print(f"le nom de fichier : {os.path.basename(doc)}")
                # Loop through all the pages
                for page_num in range(len(pdf_document)):
                    page = pdf_document.load_page(page_num)
                    images = page.get_images(full=True)
                    
                    if images:
                        text_image = get_text_from_image(images,page_num,pdf_document)
                        file_content += text_image + "\n"
                    else:
                        page_text = page.get_text()
                        if page_text:
                            file_content += page_text + "\n"
                            
                # Close the PDF file
                pdf_document.close()

            # Handle plain text (.txt) files
            elif doc.endswith(".txt"):
                with open(doc, 'r', encoding="utf-8") as f:
                    file_content += f.read()
            
            # Handle .docx (Word) files
            elif doc.endswith(".docx"):
                word_document = fitz.open(doc)

                for page_num in range(word_document.page_count):
                    page = word_document.load_page(page_num)

                    images = page.get_images(full=True)
                    if images:
                        text_image = get_text_from_image(images,page_num,word_document)
                        file_content += text_image + "\n"
                    else:
                        # If no images, read normal text
                        page_text = page.get_text()
                        if page_text:
                            file_content += page_text + "\n"
                            
                # Close the WORD file
                word_document.close()
            elif doc.endswith(".pptx"):
                # Open the PowerPoint presentation
                presentation = Presentation(doc)
                file_content = ""

                pptx_zip, image_files = extract_images_from_pptx(doc)
                
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        # Check if the shape contains text
                        if hasattr(shape, "text") and shape.text:
                            file_content += shape.text + "\n"
                        if shape.shape_type == 13:
                            for img_file in image_files:
                                image_bytes = pptx_zip.read(img_file)
                                text_image = get_text_from_image(image_bytes)
                                file_content += text_image + "\n"
                            
                pptx_zip.close()
                                
            # Handle other generic file types like .rtf or others
            else:
                with open(doc, 'r', encoding='utf-8') as file:
                    file_content += file.read()
                    
        except Exception as e:
            file_content += f"\nError reading file: {e}\n"
        
        # Append this file's content to the main text
        text += file_content + "\n" + "-" * 80 + "\n"

    return text

def get_text_from_doc(docs, descriptions=None):
    text = ""
    
    for i, doc in enumerate(docs):
        file_content = ""
        file_name = os.path.basename(doc)
        description = descriptions[i] if descriptions and i < len(descriptions) else "No description provided"
        
        file_content += f"Name of the file: {file_name}\n"
        file_content += f"Description: {description}\n"
        file_content += "Content:\n"
        
        try:
            # Handle PDF files
            if doc.endswith(".pdf"):
                pdf_document = fitz.open(doc)
                print(f"le nom de fichier : {os.path.basename(doc)}")
                # Loop through all the pages
                for page_num in range(len(pdf_document)):
                    page = pdf_document.load_page(page_num)
                    images = page.get_images(full=True)
                    
                    if images:
                        for image_index, img in enumerate(images):
                            xref = img[0]
                            base_image = pdf_document.extract_image(xref)
                            image_bytes = base_image["image"]
                            
                            # Convert image to Base64
                            base64_image = image_to_base64(image_bytes)
                            
                            # Use send_post_request to extract text from the image
                            text_image = send_post_request(base64_image)
                            file_content += text_image + "\n"
                    else:
                        page_text = page.get_text()
                        if page_text:
                            file_content += page_text + "\n"
                            
                # Close the PDF file
                pdf_document.close()

            # Handle plain text (.txt) files
            elif doc.endswith(".txt"):
                with open(doc, 'r', encoding="utf-8") as f:
                    file_content += f.read()
            
            # Handle .docx (Word) files
            elif doc.endswith(".docx"):
                word_document = fitz.open(doc)

                for page_num in range(word_document.page_count):
                    page = word_document.load_page(page_num)

                    images = page.get_images(full=True)
                    if images:
                        for image_index, img in enumerate(images):
                            xref = img[0]
                            base_image = word_document.extract_image(xref)
                            image_bytes = base_image["image"]
                            
                            # Convert image to Base64
                            base64_image = image_to_base64(image_bytes)
                            
                            # Use send_post_request to extract text from the image
                            text_image = send_post_request(base64_image)
                            file_content += text_image + "\n"
                    else:
                        # If no images, read normal text
                        page_text = page.get_text()
                        if page_text:
                            file_content += page_text + "\n"
                            
                # Close the WORD file
                word_document.close()
            
            elif doc.endswith(".pptx"):
                # Open the PowerPoint presentation
                presentation = Presentation(doc)
                
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        # Check if the shape contains text
                        if hasattr(shape, "text") and shape.text:
                            file_content += shape.text + "\n"
                        if shape.shape_type == 13:  # Assuming this is for images
                            for img_file in extract_images_from_pptx(doc):
                                # Convert image to Base64
                                base64_image = image_to_base64(img_file.read())
                                
                                # Use send_post_request to extract text from the image
                                text_image = send_post_request(base64_image)
                                file_content += text_image + "\n"
            
            # Handle other generic file types like .rtf or others
            else:
                with open(doc, 'r', encoding='utf-8') as file:
                    file_content += file.read()
                    
        except Exception as e:
            file_content += f"\nError reading file: {e}\n"
        
        # Append this file's content to the main text
        text += file_content + "\n" + "-" * 80 + "\n"

    return text

import base64

def image_to_base64(image_bytes):
    # Create a BytesIO stream to hold the image data
    buffered = BytesIO(image_bytes)
    
    # Open the image file
    with Image.open(buffered) as img:
        # Create another BytesIO stream for conversion to Base64
        output_buffer = BytesIO()
        
        # Save the image in PNG format to the buffer
        img.save(output_buffer, format="PNG")
        
        # Get the binary data from the buffer
        img_bytes = output_buffer.getvalue()
        
        # Encode the binary data to Base64
        img_base64 = base64.b64encode(img_bytes).decode("utf-8")
        
    return img_base64

import requests

def send_post_request(image_path):
    url = "http://54.216.245.50:5000/api/ocr"
    secret = "yemys44jdhKblZIezZRnC1PgcwkqkJ8H4eJYYfqRGP5QQsxb"
    key = "08xkUvyOeRSLfvW-AlzdSQ"
    
    headers = {
        "SECRET": secret,
        "KEY": key
    }
    
    # Convert image to Base64
    base64_image = image_to_base64(image_path)
    
    # Prepare the payload (body)
    payload = {
        "base_64": base64_image
    }
    
    # Send the POST request
    response = requests.post(url, headers=headers, json=payload)
    
    # Check the response
    if response.status_code == 200:
        print("Request successful:", response.json())
    else:
        print("Request failed with status code:", response.status_code)
        
def stock_in_file(question,answer,upload_json):
    file_json = session_state['file_names']
    file_path = os.path.join(upload_json, file_json)
    datafile = {'question': question, 'answer': answer}
    
    if question.strip() and answer.strip():
        if os.path.isfile(file_path):
            content_file = read_json_file(upload_json, file_json)
            content_file.append(datafile)
            with open(file_path, 'w') as file:
                json.dump(content_file, file, indent=4)
        else:
            with open(file_path, 'w') as file:
                json.dump([datafile], file, indent=4)
    else:
        with open(file_path, 'w') as file:
            json.dump([], file, indent=4)
        print("Either 'question' or 'answer' is empty. Please provide both.")

def read_json_file(directory, filename):
    file_path = os.path.join(directory, filename)
    
    if file_path:
        # Check if the file exists
        if os.path.isfile(file_path):
            print(f"Found file: {file_path}")
            # Open and read the JSON file
            with open(file_path, 'r') as file:
                data = json.load(file)
            # async with aiofiles.open(file_path, mode='r') as file:
            #     data = await file.read()
                return data
        else:
            return None
    else:
        return None

async def save_file(file,upload_folder):
    file_path = os.path.join(upload_folder, file.filename)
    try:
        async with aiofiles.open(file_path, 'wb') as f:
            await f.write(file.read())
        print(f"File saved: {file_path}")
    except Exception as e:
        print(f"Error saving file: {e}")
        return None
    return file_path

def traiter_file_have_history(history_id,history_file,upload_folder):
    file_labels = database.get_file_by_id(history_id)
    file_path = []
    descriptions = []
    
    for file_label in file_labels:
        descriptions.append(file_label[2])
        file_path.append(os.path.join(upload_folder, file_label[1]))
    
    session_state['file_names'] = history_file
    process_file(file_path,descriptions)
    
def data_history(history_file,folder_json):
    content_file = []
    if os.path.exists(os.path.join(folder_json, history_file)):
        content_file = read_json_file(folder_json, history_file)
        if content_file is None:
            return {"error": "History file not found"}
        elif content_file == []:
            return [] 
    
    return content_file

async def upload_other_files(files,upload_folder,identifiant,descriptions,user_id):
    initialize_data()
    if not files:
        return {"error": "No files provided"}
    file_stock = []
    
    file_name_json = identifiant + word_json
    
    file_paths = await asyncio.gather(*[save_file(file,upload_folder) for file in files])
    
    for file_path in file_paths:
        if file_path:
            if os.path.isfile(file_path):
                file_name = os.path.basename(file_path)
                file_stock.append(file_name)
            else:
                print(f"File does not exist: {file_path}")
        else:
            print("File path is None or empty.")
            
    session_state['hold_descriptions'] = descriptions
    
    session_state['file_names'] = file_name_json
    database.process_insert_into_files_and_history(file_stock, file_name_json,session_state['hold_descriptions'],identifiant,user_id)

    session_state['hold_files'].extend(file_paths)
    process_file(session_state['hold_files'],session_state['hold_descriptions']);

    return {"message": "Files uploaded and processed successfully"}

async def upload_file(files,upload_json,upload_folder,descriptions,user_id):
    initialize_data()
    all_filenames = []
    if not files:
        return {"error": "No files provided"}
    else: 
        print('files =>',files)
        for file in files:
            filename = file.filename
            if filename.endswith(('.pdf', '.txt', '.docx')):
                filename = filename.rsplit('.', 1)[0]
                print(filename)
                all_filenames.append(filename)
    full_filenames = ' '.join(all_filenames)   
    file_stock = []
    
    now = datetime.now()
    # random_string = ''.join(random.choices(string.ascii_letters + string.digits, k=5))
    # identifiant = 'file_'+random_string+'_'+ now.strftime("%Y-%m-%d %H:%M:%S")
    # file_name_json = 'file_'+random_string+'_'+ now.strftime("%Y-%m-%d %H-%M-%S")+word_json
    
    identifiant = full_filenames
    file_name_json = full_filenames + ' ' + now.strftime("%Y-%m-%d %H-%M-%S") + word_json
    
    file_paths = await asyncio.gather(*[save_file(file,upload_folder) for file in files])
    
    for file_path in file_paths:
        if file_path:
            if os.path.isfile(file_path):
                file_name = os.path.basename(file_path)
                file_stock.append(file_name)
            else:
                print(f"File does not exist: {file_path}")
        else:
            print("File path is None or empty.")
            
    session_state['hold_descriptions'] = descriptions
    session_state['file_names'] = file_name_json
    
    history_id = database.process_insert_into_files_and_history(file_stock, file_name_json,session_state['hold_descriptions'],identifiant,user_id)
    stock_in_file('','',upload_json)
    
    # session_state['hold_descriptions'].append(description)
    session_state['hold_files'].extend(file_paths)
    process_file(session_state['hold_files'],session_state['hold_descriptions']);

    return {"message": "Files uploaded and processed successfully", "history_id": history_id, "file_name_json": file_name_json}


def process_file(file_paths, descriptions):
    if not file_paths or not descriptions or len(file_paths) != len(descriptions):
        return {"error": "Mismatch between file paths and descriptions"}

    raw_text = get_text_from_doc(file_paths, descriptions)
    if not raw_text:
        return {"error": "No text extracted from the files"}
    
    text_size = len(raw_text)  # Use len(raw_text) for characters or tokenize it for word/token count
    if text_size < 500:
        do_split = False  # Small text
    elif 500 <= text_size <= 2000:
        do_split = True  # Medium text
    else:
        do_split = True  # Large text
   
    text_chunks = chatpdf_processor.process_text(raw_text, do_split=do_split)

    vectorstore = chatpdf_processor.get_vectorstore(text_chunks, "gpt-4-turbo")
    session_state["conversation"] = chatpdf_processor.get_conversation_chain(vectorstore)
    
        
def get_response(user_question, upload_json, emit_func):
    if not user_question:
        emit_func('error', {'message': "No question provided"})
        return

    if not session_state.get("conversation"):
        emit_func('error', {'message': "Conversation not initialized"})
        return
    
    try:
        # Initialize the streaming callback handler
        streaming_callback = CustomStreamingCallbackHandler()
        
        response = session_state["conversation"].run({"question": user_question},callbacks=[streaming_callback])
        print("\nStreaming response: ")

        full_response = streaming_callback.get_full_response()

        question = response['question']
        answer = full_response  # Use the cleaned response
        stock_in_file(question, answer, upload_json)

    except Exception as e:
        emit_func('error', {'message': str(e)})

async def upload_and_process_files(new_files,upload_folder,history_id,new_descriptions):
    initialize_data()
    history_id = history_id
    
    file_labels = database.get_file_by_id(history_id)
    old_file_paths = []
    new_file_stock = []
    old_descriptions = []
    if not new_files:
        return {"error": "No files provided"}  
    
    # Save the new files to the upload folder
    new_file_paths = await asyncio.gather(*[save_file(file, upload_folder) for file in new_files])
    
    # Process new file paths, add them to file stock
    for file_path in new_file_paths:
        if file_path:
            if os.path.isfile(file_path):
                file_name = os.path.basename(file_path)
                new_file_stock.append(file_name)
            else:
                print(f"File does not exist: {file_path}")
        else:
            print("File path is None or empty.")
            
    # Insert new files and descriptions into database
    database.process_insert_into_other_files_history(new_file_stock, history_id, new_descriptions)
    
    for file_label in file_labels:
        old_descriptions.append(file_label[2])
        old_file_paths.append(os.path.join(upload_folder, file_label[1]))
        
    file_paths = old_file_paths + new_file_paths
    descriptions = old_descriptions + new_descriptions
    session_state['hold_files'].extend(file_paths)
    session_state['hold_descriptions'].extend(descriptions)
    process_file(session_state['hold_files'], session_state['hold_descriptions'])

    return {"message": "Other files uploaded and processed successfully"}

def get_new_conversation_by_section(section_label,upload_json,upload_folder,user_id):
    initialize_data()
    if not section_label:
        return {"error": "No section provided"}
    else:
        file_labels = database.get_file_by_section(section_label)
        file_paths = []
        descriptions = []

        if file_labels:
            for file_label in file_labels:
                descriptions.append(file_label[2])
                file_paths.append(file_label[1])
             
        session_state['hold_descriptions'] = descriptions
        
        full_filenames = ' '.join(file_paths)
        now = datetime.now()
        
        identifiant = full_filenames
        file_name_json = full_filenames + ' ' + now.strftime("%Y-%m-%d %H-%M-%S") + word_json
        session_state['file_names'] = file_name_json
        
        database.process_insert_into_files_and_history(file_paths, file_name_json,session_state['hold_descriptions'],identifiant,user_id)
        stock_in_file('','',upload_json)
        
        session_state['hold_files'].extend(file_paths)
        process_file(session_state['hold_files'],session_state['hold_descriptions']);

    return {"message": "Files section uploaded and processed successfully"}

def file_exists_and_is_empty(folder_json_path, file_name):
    file_path = os.path.join(folder_json_path, file_name)
    
    # Check if file exists
    if os.path.isfile(file_path):
        # Check if file is empty (size is 0 bytes)
        return os.path.getsize(file_path) == 0
    return False

def delete_record_from_chatpdf_history_if_stay_vide(history_id,folder_json_path):
    time.sleep(900)
    file_json_name = database.get_file_json_by_id(history_id)
    file_path = os.path.join(folder_json_path, file_json_name)
    if os.path.isfile(file_path):
        if file_exists_and_is_empty(folder_json_path, file_json_name):
            print(f"File '{file_json_name}' exists and is empty.")
                # Delete the record from the database if the file is empty
            database.delete_record_from_chatpdf_history(history_id)
                # Delete the file after the database record is removed
            os.remove(file_path)
        return {"status": "deleted"}
    else:
        print(f"File '{file_json_name}' does not exist.")
        return {"status": "not exist"}
